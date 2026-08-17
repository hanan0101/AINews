# This file is part of the AI newsletter system.
"""PowerPoint export that mirrors the browser/PDF newsletter preview."""

from __future__ import annotations

import io
import re
import sys
import zipfile
from html import escape as html_escape
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from backend.utils.pdf_export_service import (
    EDGE_EXECUTABLES,
    VENV_SITE_PACKAGES,
    build_preview_pdf_document,
    harden_export_page,
    normalize_pdf_dimension,
    start_export_watchdog,
)


TEXT_SELECTOR = ", ".join(
    (
        ".page-title h1",
        ".section-divider h2",
        ".news-card .card-title",
        ".news-card .card-text",
        ".news-card .card-source",
        ".feature-card:not(.movie-card) .card-title",
        ".feature-card:not(.movie-card) .card-text",
        ".feature-card:not(.movie-card) .feature-link",
        ".feature-card:not(.movie-card) .feature-meta > span:not(.level-dots)",
        ".movie-card .movie-title",
        ".movie-card .movie-text",
        ".cultural-assistant-copy",
        ".cultural-assistant-start",
        ".page-footer",
    )
)
IMAGE_SELECTOR = ", ".join(
    (
        ".obs-logo",
        ".moc-logo",
        ".card-logo img",
        ".movie-poster img",
        ".cultural-assistant-robot",
    )
)
FONT_FAMILY = "Effra"
PROJECT_FONT_FILES = (
    ("regular", "Effra_Rg.ttf"),
    ("bold", "Effra_Md.ttf"),
)


def _rgb(css_color: str, fallback=(46, 36, 29)) -> RGBColor:
    values = re.findall(r"[\d.]+", str(css_color or ""))
    if len(values) >= 3:
        return RGBColor(*(max(0, min(255, round(float(value)))) for value in values[:3]))
    return RGBColor(*fallback)


def _set_rtl(paragraph):
    paragraph.alignment = PP_ALIGN.RIGHT
    p_pr = paragraph._p.get_or_add_pPr()
    if p_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rtl") is None:
        rtl = OxmlElement("a:rtl")
        rtl.set("val", "1")
        p_pr.append(rtl)


def _powerpoint_segments(value, default_direction="rtl") -> list[tuple[str, str]]:
    """Convert browser bidi isolates into native directional PowerPoint runs."""
    directions = [default_direction]
    buffer = []
    segments = []

    def flush():
        if buffer:
            segments.append(("".join(buffer), directions[-1]))
            buffer.clear()

    for character in str(value or ""):
        if character in {"\u2066", "\u2068"}:
            flush()
            directions.append("ltr")
        elif character == "\u2067":
            flush()
            directions.append("rtl")
        elif character == "\u2069":
            flush()
            if len(directions) > 1:
                directions.pop()
        elif character in {"\u202a", "\u202b", "\u202c", "\u202d", "\u202e"}:
            continue
        else:
            buffer.append(character)
    flush()
    expanded = []
    latin_chunk = re.compile(
        r"(\([^()\u0600-\u06ff]*[A-Za-z0-9][^()\u0600-\u06ff]*\)"
        r"|[A-Za-z0-9][A-Za-z0-9._+\-/]*(?: [A-Za-z0-9][A-Za-z0-9._+\-/]*)*)"
    )
    for text, direction in segments:
        if direction != "rtl":
            expanded.append((text, direction))
            continue
        for part in latin_chunk.split(text):
            if not part:
                continue
            part_direction = "ltr" if latin_chunk.fullmatch(part) else "rtl"
            expanded.append((part, part_direction))
    return expanded


def _set_powerpoint_run_font(run, *, size, color, bold, direction="rtl"):
    """Set Latin and complex-script font data so Arabic does not fall back."""
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.underline = False
    run.font.color.rgb = color
    run_properties = run._r.get_or_add_rPr()
    run_properties.set("lang", "en-US" if direction == "ltr" else "ar-SA")
    run_properties.set("altLang", "ar-SA" if direction == "ltr" else "en-US")
    run_properties.set("rtl", "1" if direction == "rtl" else "0")
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = run_properties.find(
            f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':', 1)[1]}"
        )
        if node is None:
            node = OxmlElement(tag)
            run_properties.append(node)
        node.set("typeface", FONT_FAMILY)


def _embed_project_fonts(pptx_bytes: bytes) -> bytes:
    """Embed the two editable Effra faces used by the browser/PDF preview."""
    font_root = Path(__file__).resolve().parents[2] / "frontend" / "fonts"
    fonts = []
    for style, filename in PROJECT_FONT_FILES:
        font_path = font_root / filename
        if not font_path.is_file():
            raise RuntimeError(f"PowerPoint font file is missing: {font_path}")
        fonts.append((style, filename, font_path.read_bytes()))

    presentation_part = "ppt/presentation.xml"
    relationships_part = "ppt/_rels/presentation.xml.rels"
    content_types_part = "[Content_Types].xml"
    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as source:
        files = {name: source.read(name) for name in source.namelist()}

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ET.register_namespace("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
    ET.register_namespace("p", p_ns)
    ET.register_namespace("r", r_ns)
    ET.register_namespace("", rel_ns)

    rel_root = ET.fromstring(files[relationships_part])
    used_ids = {
        relationship.get("Id", "")
        for relationship in rel_root.findall(f"{{{rel_ns}}}Relationship")
    }

    def next_relationship_id():
        candidate = 1
        while f"rId{candidate}" in used_ids:
            candidate += 1
        value = f"rId{candidate}"
        used_ids.add(value)
        return value

    font_relationships = []
    for index, (style, filename, font_bytes) in enumerate(fonts, start=1):
        relationship_id = next_relationship_id()
        part_name = f"ppt/fonts/font{index}.fntdata"
        relationship = ET.SubElement(rel_root, f"{{{rel_ns}}}Relationship")
        relationship.set("Id", relationship_id)
        relationship.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font",
        )
        relationship.set("Target", f"fonts/font{index}.fntdata")
        files[part_name] = font_bytes
        font_relationships.append((style, relationship_id, part_name))
    files[relationships_part] = ET.tostring(
        rel_root, encoding="utf-8", xml_declaration=True
    )

    presentation_root = ET.fromstring(files[presentation_part])
    existing = presentation_root.find(f"{{{p_ns}}}embeddedFontLst")
    if existing is not None:
        presentation_root.remove(existing)
    embedded_list = ET.Element(f"{{{p_ns}}}embeddedFontLst")
    embedded_font = ET.SubElement(embedded_list, f"{{{p_ns}}}embeddedFont")
    font = ET.SubElement(embedded_font, f"{{{p_ns}}}font")
    font.set("typeface", FONT_FAMILY)
    font.set("pitchFamily", "34")
    font.set("charset", "0")
    for style, relationship_id, _part_name in font_relationships:
        face = ET.SubElement(embedded_font, f"{{{p_ns}}}{style}")
        face.set(f"{{{r_ns}}}id", relationship_id)
    notes_size = presentation_root.find(f"{{{p_ns}}}notesSz")
    insert_at = (
        list(presentation_root).index(notes_size) + 1
        if notes_size is not None
        else len(presentation_root)
    )
    presentation_root.insert(insert_at, embedded_list)
    files[presentation_part] = ET.tostring(
        presentation_root, encoding="utf-8", xml_declaration=True
    )

    ET.register_namespace("", ct_ns)
    content_types_root = ET.fromstring(files[content_types_part])
    existing_parts = {
        override.get("PartName")
        for override in content_types_root.findall(f"{{{ct_ns}}}Override")
    }
    for _family, _relationship_id, part_name in font_relationships:
        package_name = f"/{part_name}"
        if package_name in existing_parts:
            continue
        override = ET.SubElement(content_types_root, f"{{{ct_ns}}}Override")
        override.set("PartName", package_name)
        override.set("ContentType", "application/x-fontdata")
    files[content_types_part] = ET.tostring(
        content_types_root, encoding="utf-8", xml_declaration=True
    )

    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as target:
        for name, data in files.items():
            target.writestr(name, data)
    return output.getvalue()


def _render_preview_layers(
    preview_html: str,
    width: float,
    height: float,
    origin: str,
    direction: str,
    source_file: str,
) -> dict:
    """Render an exact decoration layer and extract editable overlay geometry."""
    width = normalize_pdf_dimension(width, 1000)
    height = normalize_pdf_dimension(height, 1340)
    document_html = build_preview_pdf_document(
        preview_html,
        width,
        height,
        origin,
        direction,
        scale=1,
        source_file=source_file,
    )
    if VENV_SITE_PACKAGES.exists() and str(VENV_SITE_PACKAGES) not in sys.path:
        sys.path.insert(0, str(VENV_SITE_PACKAGES))
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise RuntimeError(
            "Playwright is required for PowerPoint export. "
            "Run: pip install playwright && playwright install chromium"
        ) from exc

    with sync_playwright() as playwright:
        edge_path = next((path for path in EDGE_EXECUTABLES if path.exists()), None)
        try:
            browser = playwright.chromium.launch(
                executable_path=str(edge_path) if edge_path else None
            )
        except Exception as exc:
            raise RuntimeError(
                "Chromium/Edge could not start for PowerPoint export. "
                "Install Microsoft Edge or run: playwright install chromium"
            ) from exc
        watchdog = start_export_watchdog(browser)
        try:
            page = browser.new_page(
                viewport={"width": int(width), "height": int(height)},
                device_scale_factor=2,
                java_script_enabled=False,
            )
            harden_export_page(page, origin)
            page.emulate_media(media="screen")
            page.set_content(document_html, wait_until="networkidle", timeout=30000)
            page.evaluate(
                """async () => {
                  if (document.fonts && document.fonts.ready) await document.fonts.ready;
                  await Promise.all(Array.from(document.images || []).map(img =>
                    img.complete ? true : new Promise(resolve => {
                      img.onload = resolve;
                      img.onerror = resolve;
                      setTimeout(resolve, 12000);
                    })
                  ));
                }"""
            )

            text_layers = page.eval_on_selector_all(
                TEXT_SELECTOR,
                """elements => {
                  const pageRect = document.querySelector('.page').getBoundingClientRect();
                  return elements.map(element => {
                    const rect = element.getBoundingClientRect();
                    const style = getComputedStyle(element);
                    const anchor = element.closest('a[href]') || (element.matches('a[href]') ? element : null);
                    const role =
                      element.matches('.page-title h1') ? 'page-title' :
                      element.matches('.section-divider h2') ? 'section-title' :
                      element.matches('.news-card .card-title') ? 'news-title' :
                      element.matches('.news-card .card-text') ? 'news-body' :
                      element.matches('.news-card .card-source') ? 'news-source' :
                      element.matches('.feature-card:not(.movie-card) .card-title') ? 'course-title' :
                      element.matches('.feature-card:not(.movie-card) .card-text') ? 'course-body' :
                      element.matches('.feature-card:not(.movie-card) .feature-link') ? 'course-link' :
                      element.matches('.feature-card:not(.movie-card) .feature-meta > span:not(.level-dots)') ? 'course-level' :
                      element.matches('.movie-title') ? 'movie-title' :
                      element.matches('.movie-text') ? 'movie-body' :
                      element.matches('.cultural-assistant-copy') ? 'assistant-copy' :
                      element.matches('.cultural-assistant-start') ? 'assistant-start' :
                      element.matches('.page-footer') ? 'footer' : 'text';
                    const rawText = String(element.innerText || element.textContent || '');
                    return {
                      text: rawText
                        .replace(/[\\t\\f\\v ]+/g, ' ')
                        .replace(/ *\\n+ */g, '\\n')
                        .trim(),
                      role,
                      x: rect.left - pageRect.left,
                      y: rect.top - pageRect.top,
                      w: rect.width,
                      h: rect.height,
                      color: style.color,
                      fontFamily: style.fontFamily,
                      fontSize: parseFloat(style.fontSize) || 14,
                      fontWeight: parseInt(style.fontWeight, 10) || 400,
                      textAlign: style.textAlign,
                      direction: style.direction,
                      lineHeight: parseFloat(style.lineHeight) || 0,
                      href: anchor ? anchor.href : ''
                    };
                  }).filter(item => item.text && item.w > 1 && item.h > 1);
                }""",
            )

            image_descriptors = page.eval_on_selector_all(
                IMAGE_SELECTOR,
                """elements => {
                  const pageRect = document.querySelector('.page').getBoundingClientRect();
                  const intersect = (a, b) => ({
                    left: Math.max(a.left, b.left),
                    top: Math.max(a.top, b.top),
                    right: Math.min(a.right, b.right),
                    bottom: Math.min(a.bottom, b.bottom)
                  });
                  return elements.map(element => {
                    const rect = element.getBoundingClientRect();
                    let visible = {
                      left: rect.left, top: rect.top,
                      right: rect.right, bottom: rect.bottom
                    };
                    let ancestor = element.parentElement;
                    while (ancestor) {
                      const style = getComputedStyle(ancestor);
                      const clipsX = ['hidden', 'clip', 'scroll', 'auto'].includes(style.overflowX);
                      const clipsY = ['hidden', 'clip', 'scroll', 'auto'].includes(style.overflowY);
                      if (clipsX || clipsY) {
                        const clip = ancestor.getBoundingClientRect();
                        if (clipsX) {
                          visible.left = Math.max(visible.left, clip.left);
                          visible.right = Math.min(visible.right, clip.right);
                        }
                        if (clipsY) {
                          visible.top = Math.max(visible.top, clip.top);
                          visible.bottom = Math.min(visible.bottom, clip.bottom);
                        }
                      }
                      ancestor = ancestor.parentElement;
                    }
                    visible = intersect(visible, pageRect);
                    const style = getComputedStyle(element);
                    return {
                      src: element.currentSrc || element.src || '',
                      x: visible.left - pageRect.left,
                      y: visible.top - pageRect.top,
                      w: Math.max(0, visible.right - visible.left),
                      h: Math.max(0, visible.bottom - visible.top),
                      fullW: rect.width,
                      fullH: rect.height,
                      offsetX: rect.left - visible.left,
                      offsetY: rect.top - visible.top,
                      objectFit: style.objectFit || 'contain',
                      objectPosition: style.objectPosition || '50% 50%'
                    };
                  }).filter(item => item.src && item.w > 1 && item.h > 1);
                }""",
            )
            image_layers = []
            capture_page = browser.new_page(
                viewport={"width": 1000, "height": 1000},
                device_scale_factor=2,
                java_script_enabled=False,
            )
            harden_export_page(capture_page, origin)
            try:
                for descriptor in image_descriptors:
                    capture_page.set_viewport_size(
                        {
                            "width": max(1, round(descriptor["w"])),
                            "height": max(1, round(descriptor["h"])),
                        }
                    )
                    capture_html = f"""<!doctype html><html><head><base href="{html_escape(origin, quote=True)}/">
<style>
html,body{{margin:0;padding:0;background:transparent;overflow:hidden}}
#clip{{position:relative;width:{descriptor["w"]}px;height:{descriptor["h"]}px;overflow:hidden;background:transparent}}
#clip img{{position:absolute;left:{descriptor["offsetX"]}px;top:{descriptor["offsetY"]}px;
width:{descriptor["fullW"]}px;height:{descriptor["fullH"]}px;
object-fit:{html_escape(descriptor["objectFit"])};object-position:{html_escape(descriptor["objectPosition"])}}}
</style></head><body><div id="clip"><img src="{html_escape(descriptor["src"], quote=True)}"></div></body></html>"""
                    capture_page.set_content(capture_html, wait_until="networkidle", timeout=30000)
                    capture_page.evaluate(
                        """async () => {
                          const image = document.querySelector('img');
                          if (!image.complete) await new Promise(resolve => {
                            image.onload = resolve; image.onerror = resolve;
                            setTimeout(resolve, 12000);
                          });
                        }"""
                    )
                    image_layers.append(
                        {
                            "x": descriptor["x"],
                            "y": descriptor["y"],
                            "w": descriptor["w"],
                            "h": descriptor["h"],
                            "bytes": capture_page.locator("#clip").screenshot(
                                type="png", omit_background=True
                            ),
                        }
                    )
            finally:
                capture_page.close()

            page.evaluate(
                """({textSelector, imageSelector}) => {
                  const walkerFor = element => {
                    const walker = document.createTreeWalker(element, NodeFilter.SHOW_TEXT);
                    const nodes = [];
                    while (walker.nextNode()) nodes.push(walker.currentNode);
                    nodes.forEach(node => {
                      const span = document.createElement('span');
                      span.style.color = 'transparent';
                      span.style.textShadow = 'none';
                      node.parentNode.insertBefore(span, node);
                      span.appendChild(node);
                    });
                  };
                  document.querySelectorAll(textSelector).forEach(walkerFor);
                  document.querySelectorAll(imageSelector).forEach(element => {
                    element.style.visibility = 'hidden';
                  });
                }""",
                {"textSelector": TEXT_SELECTOR, "imageSelector": IMAGE_SELECTOR},
            )
            background = page.locator(".page").screenshot(type="png")
            return {
                "width": width,
                "height": height,
                "background": background,
                "texts": text_layers,
                "images": image_layers,
            }
        finally:
            watchdog.cancel()
            browser.close()


def _presentation_from_layers(layers: dict, title: str = "") -> bytes:
    """Create one editable PowerPoint page from rendered browser layers."""
    pixel_width = float(layers["width"])
    pixel_height = float(layers["height"])
    slide_width_in = 10.0
    slide_height_in = slide_width_in * pixel_height / pixel_width
    scale_in = slide_width_in / pixel_width

    presentation = Presentation()
    presentation.slide_width = Inches(slide_width_in)
    presentation.slide_height = Inches(slide_height_in)
    presentation.core_properties.title = str(title or "AI Newsletter")
    presentation.core_properties.subject = "Editable newsletter matching the PDF preview"
    presentation.core_properties.author = "AI Newsletter System"
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background_shape = slide.shapes.add_picture(
        io.BytesIO(layers["background"]),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    background_shape.name = "Newsletter design (matches PDF)"

    for index, item in enumerate(layers.get("texts") or [], start=1):
        role = str(item.get("role") or "text")
        x = float(item["x"])
        y = float(item["y"])
        w = max(1.0, float(item["w"]))
        h = max(1.0, float(item["h"]))
        if role == "section-title":
            # PowerPoint places the Effra baseline lower than Chromium. Lift
            # the editable heading so the divider stays centered beside it.
            x -= 10
            y -= 5
            w += 20
            h += 5
        elif role == "course-title":
            y -= 1.5
            h += 5
        elif role == "course-body":
            y -= 4.5
            h += 6
        elif role in {"course-link", "course-level"}:
            y -= 8
            h += 10
            if role == "course-level":
                x -= 55
                w += 55
        left = Inches(x * scale_in)
        top = Inches(y * scale_in)
        width = Inches(w * scale_in)
        height = Inches(h * scale_in)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.name = f"Editable newsletter text {index:02d}"
        if role == "section-title":
            # The divider is baked into the exact background image. An opaque
            # paper-coloured fill prevents it from showing through editable
            # PowerPoint text while preserving both side segments.
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(245, 237, 223)
        else:
            box.fill.background()
        box.line.fill.background()
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        compact_line = h <= max(32, float(item.get("fontSize") or 14) * 1.8)
        top_anchored_roles = {
            "page-title",
            "section-title",
            "news-title",
            "news-body",
            "course-title",
            "course-body",
            "movie-title",
            "movie-body",
            "assistant-copy",
            "footer",
            "course-link",
            "course-level",
        }
        frame.vertical_anchor = (
            MSO_ANCHOR.TOP
            if role in top_anchored_roles
            else MSO_ANCHOR.MIDDLE if compact_line else MSO_ANCHOR.TOP
        )
        paragraph = frame.paragraphs[0]
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        alignment = str(item.get("textAlign") or "").lower()
        paragraph.alignment = {
            "center": PP_ALIGN.CENTER,
            "left": PP_ALIGN.LEFT,
            "right": PP_ALIGN.RIGHT,
        }.get(alignment, PP_ALIGN.RIGHT)
        raw_text = str(item.get("text") or "")
        contains_arabic = bool(re.search(r"[\u0600-\u06ff]", raw_text))
        text_direction = str(item.get("direction") or "").lower()
        paragraph_direction = "rtl" if text_direction == "rtl" and contains_arabic else "ltr"
        if paragraph_direction == "rtl":
            _set_rtl(paragraph)
        href = str(item.get("href") or "")
        font_scale = {
            "course-title": 1.04,
            "course-body": 1.03,
            "course-link": 1.03,
            "course-level": 1.03,
        }.get(role, 1.0)
        font_size = max(
            6,
            float(item.get("fontSize") or 14) * scale_in * 72 * font_scale,
        )
        css_line_height = float(item.get("lineHeight") or 0)
        if css_line_height > 0 and role not in {"page-title", "section-title"}:
            line_height_scale = 0.76 if role == "course-body" else 1.0
            paragraph.line_spacing = Pt(
                css_line_height * scale_in * 72 * line_height_scale
            )
        for segment_text, segment_direction in _powerpoint_segments(
            raw_text, paragraph_direction
        ):
            run = paragraph.add_run()
            run.text = segment_text
            _set_powerpoint_run_font(
                run,
                size=font_size,
                color=_rgb(item.get("color")),
                bold=(
                    int(item.get("fontWeight") or 400) >= 600
                    or role == "course-title"
                ),
                direction=segment_direction,
            )
        # A run-level hyperlink forces Office's blue/underline theme. Linking
        # the text box keeps the source clickable while preserving PDF styling.
        if href.startswith(("http://", "https://")):
            box.click_action.hyperlink.address = href

    for index, item in enumerate(layers.get("images") or [], start=1):
        picture = slide.shapes.add_picture(
            io.BytesIO(item["bytes"]),
            Inches(float(item["x"]) * scale_in),
            Inches(float(item["y"]) * scale_in),
            width=Inches(float(item["w"]) * scale_in),
            height=Inches(float(item["h"]) * scale_in),
        )
        picture.name = f"Editable newsletter logo/image {index:02d}"

    stream = io.BytesIO()
    presentation.save(stream)
    return _embed_project_fonts(stream.getvalue())


def export_newsletter_pptx_bytes(payload, origin=""):
    """Return an editable PPTX matching the same live preview used by PDF."""
    payload = payload if isinstance(payload, dict) else {}
    preview_html = str(payload.get("html") or "")
    if not preview_html or "page" not in preview_html:
        raise ValueError("PowerPoint export requires the current newsletter preview HTML")
    width = payload.get("width") or 1000
    height = payload.get("height") or 1340
    direction = str(payload.get("direction") or "rtl")
    source_file = str(payload.get("source_html") or "News.html")
    layers = _render_preview_layers(
        preview_html,
        width,
        height,
        origin,
        direction,
        source_file,
    )
    newsletter = payload.get("newsletter_json")
    template = newsletter.get("template") if isinstance(newsletter, dict) else {}
    title = template.get("newsletter_title") if isinstance(template, dict) else ""
    return _presentation_from_layers(layers, title)


__all__ = ["export_newsletter_pptx_bytes"]
