import io
import unittest
import zipfile

from pptx import Presentation

from backend.utils.pptx_export_service import (
    _powerpoint_segments,
    _presentation_from_layers,
)


class PowerPointExportTests(unittest.TestCase):
    def test_mixed_arabic_english_text_uses_directional_runs_without_controls(self):
        segments = _powerpoint_segments(
            "أطلقت شركة \u2066(Meta)\u2069 نموذج (Muse Spark 1.1) للمهام"
        )
        joined = "".join(text for text, _ in segments)
        self.assertEqual(
            joined,
            "أطلقت شركة (Meta) نموذج (Muse Spark 1.1) للمهام",
        )
        self.assertNotRegex(joined, r"[\u202a-\u202e\u2066-\u2069]")
        self.assertIn(("(Meta)", "ltr"), segments)
        self.assertIn(("(Muse Spark 1.1)", "ltr"), segments)

    def test_export_uses_exact_background_and_editable_overlays(self):
        layers = {
            "width": 1000,
            "height": 1500,
            "background": (
                b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
                b"\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
                b"\x00\x00\x00\rIDAT\x08\xd7c\xf8\xcf\xc0\xf0\x1f\x00"
                b"\x05\x00\x01\xff\x89\x99=\x1d\x00\x00\x00\x00IEND\xaeB`\x82"
            ),
            "texts": [
                {
                    "text": "عنوان قابل للتعديل\nالسطر الثاني",
                    "role": "page-title",
                    "x": 100,
                    "y": 120,
                    "w": 800,
                    "h": 60,
                    "fontSize": 28,
                    "fontWeight": 700,
                    "fontFamily": "Avenir",
                    "color": "rgb(46, 36, 29)",
                    "textAlign": "right",
                    "direction": "rtl",
                    "href": "https://example.com/source",
                }
            ],
            "images": [],
        }

        pptx_bytes = _presentation_from_layers(layers, "نشرة تجريبية")
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        presentation = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(len(presentation.slides), 1)
        self.assertAlmostEqual(
            presentation.slide_height / presentation.slide_width,
            1.5,
            places=2,
        )
        editable_text = [
            shape.text
            for shape in presentation.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
        ]
        self.assertIn("عنوان قابل للتعديل\nالسطر الثاني", editable_text)
        self.assertGreaterEqual(len(presentation.slides[0].shapes), 2)
        linked_shapes = [
            shape
            for shape in presentation.slides[0].shapes
            if shape.click_action.hyperlink.address
        ]
        self.assertEqual(len(linked_shapes), 1)
        self.assertEqual(
            linked_shapes[0].click_action.hyperlink.address,
            "https://example.com/source",
        )
        text_runs = linked_shapes[0].text_frame.paragraphs[0].runs
        self.assertTrue(text_runs)
        self.assertTrue(all(run.font.name == "Effra" for run in text_runs))
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as package:
            names = set(package.namelist())
            self.assertIn("ppt/fonts/font1.fntdata", names)
            self.assertIn("ppt/fonts/font2.fntdata", names)
            presentation_xml = package.read("ppt/presentation.xml")
            self.assertIn(b"embeddedFontLst", presentation_xml)
            self.assertIn(b'typeface="Effra"', presentation_xml)


if __name__ == "__main__":
    unittest.main()
